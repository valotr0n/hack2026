from __future__ import annotations

import io
import json
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from ..llm import chat

router = APIRouter()


class PresentationRequest(BaseModel):
    text: str
    title: str = ""
    style: str = "business"  # business | academic | popular
    prompt: str = ""  # дополнительные инструкции от пользователя


class PresentationResponse(BaseModel):
    slides: list[dict]  # структура для превью в браузере


# ── JSON структура слайдов ────────────────────────────────────────────────────

async def _generate_slides(text: str, title: str, style: str, prompt: str = "") -> list[dict]:
    style_prompt = {
        "business": "деловой корпоративный стиль, чёткие тезисы, минимум воды",
        "academic": "академический стиль, подробные объяснения, термины",
        "popular": "простой понятный язык, примеры, аналогии",
    }.get(style, "деловой стиль")

    system = (
        "Ты — эксперт по созданию презентаций. "
        "Используй предоставленный текст как основу, дополняя своими знаниями там где это уместно. "
        "Структурируй материал в чёткие, ёмкие слайды. "
        "Отвечай строго в формате JSON без лишнего текста."
    )

    user = (
        f"Создай презентацию в {style_prompt} по тексту ниже.\n"
        + (f'Заголовок презентации: "{title}"\n' if title else "")
        + (f"Дополнительные инструкции от пользователя: {prompt}\n" if prompt.strip() else "")
        + "Верни JSON строго в этом формате:\n"
        '{"slides": [\n'
        '  {"type": "title", "title": "Название презентации", "subtitle": "Подзаголовок"},\n'
        '  {"type": "content", "title": "Заголовок слайда", "body": "1-2 предложения объясняющие суть этого раздела", "bullets": ["тезис 1", "тезис 2", "тезис 3"]},\n'
        '  {"type": "content", "title": "Следующий раздел", "body": "1-2 предложения контекста и пояснения", "bullets": ["тезис 1", "тезис 2"]},\n'
        '  {"type": "summary", "title": "Итоги", "body": "Общий вывод 1-2 предложения", "bullets": ["вывод 1", "вывод 2"]}\n'
        "]}\n\n"
        "Обязательные правила:\n"
        "- Первый слайд type=title, последний type=summary\n"
        "- 5–8 слайдов итого\n"
        "- КАЖДЫЙ слайд кроме title ОБЯЗАН содержать поле body — 1-2 содержательных предложения\n"
        "- body пишется связным текстом, раскрывает суть слайда\n"
        "- bullets — 3–5 коротких тезисов дополняющих body\n\n"
        f"Текст:\n{text}"
    )

    raw = await chat(system=system, user=user)

    try:
        start = raw.find("{")
        end = raw.rfind("}") + 1
        data = json.loads(raw[start:end])
        return data.get("slides", [])
    except Exception:
        raise HTTPException(status_code=500, detail="Не удалось разобрать структуру презентации")


# ── Сборка PPTX ───────────────────────────────────────────────────────────────

def _build_pptx(slides: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    FONT = "Geist"

    # Чёрно-белая тема
    COLOR_BLACK = RGBColor(0x0A, 0x0A, 0x0A)
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_GRAY  = RGBColor(0x55, 0x55, 0x55)
    COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

    def _set_bg(slide, color: RGBColor) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(slide, text: str, left, top, width, height,
                     font_size: int, bold: bool = False,
                     color: RGBColor = COLOR_BLACK,
                     align=PP_ALIGN.LEFT,
                     italic: bool = False) -> None:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        title_text = slide_data.get("title", "")
        body_text = slide_data.get("body", "")
        bullets = slide_data.get("bullets", [])

        if slide_type == "title":
            layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(layout)
            _set_bg(slide, COLOR_BLACK)

            # Белая горизонтальная полоса снизу
            bar = slide.shapes.add_shape(
                1,
                Inches(0), Inches(6.8),
                Inches(13.33), Inches(0.05),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_WHITE
            bar.line.fill.background()

            _add_textbox(slide, title_text,
                         Inches(1.0), Inches(2.0), Inches(11.0), Inches(2.0),
                         font_size=44, bold=True, color=COLOR_WHITE,
                         align=PP_ALIGN.LEFT)

            subtitle = slide_data.get("subtitle", "")
            if subtitle:
                _add_textbox(slide, subtitle,
                             Inches(1.0), Inches(4.2), Inches(11.0), Inches(0.8),
                             font_size=20, color=COLOR_GRAY,
                             align=PP_ALIGN.LEFT)

        else:  # content / summary
            layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(layout)
            _set_bg(slide, COLOR_WHITE)

            # Чёрная заголовочная полоса
            header = slide.shapes.add_shape(
                1,
                Inches(0), Inches(0),
                Inches(13.33), Inches(1.2),
            )
            header.fill.solid()
            header.fill.fore_color.rgb = COLOR_BLACK
            header.line.fill.background()

            _add_textbox(slide, title_text,
                         Inches(0.4), Inches(0.1), Inches(12.5), Inches(1.0),
                         font_size=26, bold=True, color=COLOR_WHITE)

            y_offset = Inches(1.35)

            # Body текст под заголовком
            if body_text:
                _add_textbox(slide, body_text,
                             Inches(0.5), y_offset, Inches(12.0), Inches(0.9),
                             font_size=14, color=COLOR_GRAY, italic=True)
                y_offset += Inches(1.0)

            # Тонкий разделитель
            if body_text:
                sep = slide.shapes.add_shape(
                    1,
                    Inches(0.5), y_offset,
                    Inches(12.33), Inches(0.02),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                sep.line.fill.background()
                y_offset += Inches(0.15)

            # Буллеты
            bullet_text = "\n".join(f"— {b}" for b in bullets)
            remaining = Inches(7.5) - y_offset - Inches(0.2)
            _add_textbox(slide, bullet_text,
                         Inches(0.5), y_offset, Inches(12.0), remaining,
                         font_size=18, color=COLOR_BLACK)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Эндпоинты ─────────────────────────────────────────────────────────────────

@router.post(
    "/presentation/preview",
    response_model=PresentationResponse,
    summary="Структура презентации (превью)",
)
async def presentation_preview(req: PresentationRequest) -> PresentationResponse:
    slides = await _generate_slides(req.text, req.title, req.style, req.prompt)
    return PresentationResponse(slides=slides)


@router.post(
    "/presentation/download",
    summary="Скачать презентацию (PPTX)",
)
async def presentation_download(req: PresentationRequest) -> StreamingResponse:
    slides = await _generate_slides(req.text, req.title, req.style, req.prompt)
    pptx_bytes = _build_pptx(slides)
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="presentation.pptx"'},
    )
